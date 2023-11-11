from pptx import Presentation
from pptx.util import Pt, Inches

# יצירת מצגת חדשה
prs = Presentation()

# פרטי השקפים והתוכן
slides_content = [
    ("הזכות לכבוד",
     "הזכות לכבוד – זכותו של כל אדם לא להיות מושפל, שלא יפגעו בכבודו, שלא יעליבו אותו ויבזו אותו. מזכות זו נגזרות שתי זכויות."),
    ("הזכות לפרטיות", "הזכות לחיות בלי התערבות וחדירה לחיים, לגוף ולרכוש של האדם."),
    ("הזכות לשם טוב", "הזכות שלא יפרסמו על אדם מידע שקרי ושלילי."),
    ("באחריות מי לדאוג לקיום הזכות?", "תפקידה של המדינה.")
]

# הוספת השקפים והרקעים
for index, (title, content) in enumerate(slides_content):
    slide_layout = prs.slide_layouts[5]  # בחירת עיצוב השקף עם רקע ריק
    slide = prs.slides.add_slide(slide_layout)

    # הוספת רקע גרפי
    # יש להחליף את 'background.png' עם הנתיב לתמונת הרקע שברצונך להשתמש בה
    bg_image_path = f'background_{index + 1}.png'  # נניח שיש לך תמונות רקע שונות לכל שקף
    slide.shapes.add_picture(bg_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # הוספת כותרת ותוכן על גבי הרקע
    title_shape = slide.shapes.title
    title_shape.text = title
    content_shape = slide.placeholders[0]
    content_shape.text = content

    # עיצוב הטקסט כאן (כמו בדוגמה הקודמת)

# שמירת המצגת
prs.save('presentation_with_backgrounds.pptx')
